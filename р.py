import requests
import telebot
import json
import os
import re
from pptx import Presentation


bot = telebot.TeleBot(TOKEN)


def ask_yandex_gpt(messages):
    url = "https://llm.api.cloud.yandex.net/foundationModels/v1/completion"

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Api-Key {API_KEY}"
    }

    data = {
        "modelUri": f"gpt://{FOLDER_ID}/yandexgpt",
        "completionOptions": {
            "stream": False,
            "temperature": 0.6,
            "maxTokens": 2000
        },
        "messages": messages
    }

    try:
        response = requests.post(url, headers=headers, json=data, timeout=60)
        response.raise_for_status()
        result = response.json()
        return result["result"]["alternatives"][0]["message"]["text"]
    except Exception as e:
        print("Ошибка YandexGPT:", e)
        return None


# джсонджсонджсон
def extract_json(text):
    try:
        return json.loads(text)
    except:
        match = re.search(r'\{.*\}', text, re.DOTALL)
        if match:
            return json.loads(match.group())
        else:
            raise ValueError("JSON не найден")


# создание пптх
def create_presentation(slides_data, filename):
    prs = Presentation()

    for i, slide_data in enumerate(slides_data):
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = slide_data.get("title", "Без названия")

        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_data.get("content", "")

    prs.save(filename)


@bot.message_handler(commands=['start'])
def start_command(message):
    user_id = message.chat.id
    help_text = '''Привет! Я бот-учитель. Чтобы начать, отправь мне сообщение в формате:

    Название предмета
    Область знаний
    Длительность урока (в минутах)

    Например:
    История
    Древний мир
    40

    Используй /clear чтобы очистить историю диалога.
    Используй /help чтобы увидеть это сообщение снова.'''
    bot.send_message(user_id, help_text)


@bot.message_handler(func=lambda message: True, content_types=['text'])
def handle_message(message):
    user_id = message.chat.id
    user_text = message.text.strip()

    if user_text.startswith('/'):
        return

    lines = user_text.split('\n')
    subject = lines[0] if len(lines) > 0 else "предмет"
    area = lines[1] if len(lines) > 1 else "общая область"
    duration = lines[2] if len(lines) > 2 else "45"

    system_message = f"""
Ты профессиональный учитель {subject}, эксперт в области {area}.

Создай:
1) Подробный сценарий урока на {duration} минут
2) Структуру презентации

Ответ верни в JSON формате:

{{
  "lesson_plan": "Текст сценария...",
  "slides": [
    {{
      "title": "Заголовок",
      "content": "Текст через \\n"
    }}
  ]
}}
"""

    messages = [
        {"role": "system", "text": system_message},
        {"role": "user", "text": user_text}
    ]

    bot.send_message(user_id, "Генерирую урок и презентацию...")

    answer = ask_yandex_gpt(messages)

    if not answer:
        bot.send_message(user_id, "Ошибка генерации.")
        return

    try:
        data = extract_json(answer)

        lesson_plan = data.get("lesson_plan", "Сценарий не найден.")
        slides = data.get("slides", [])

        for i in range(0, len(lesson_plan), 4000):
            bot.send_message(user_id, lesson_plan[i:i + 4000])

        # презу создаем
        filename = f"presentation_{user_id}.pptx"
        create_presentation(slides, filename)

        with open(filename, "rb") as file:
            bot.send_document(user_id, file)

        os.remove(filename)

    except Exception as e:
        print("Ошибка обработки:", e)
        bot.send_message(user_id, "Ошибка при создании презентации.")


if __name__ == "__main__":
    print("Бот запущен...")
    bot.polling(none_stop=True)